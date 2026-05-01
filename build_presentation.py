"""
Run:  python3 build_presentation.py
Output: vortex_presentation.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
DARK_BLUE  = RGBColor(0x0D, 0x2B, 0x55)
MID_BLUE   = RGBColor(0x1A, 0x5C, 0x9E)
LIGHT_BLUE = RGBColor(0xD6, 0xE8, 0xF7)
ACCENT     = RGBColor(0xE8, 0x4A, 0x1A)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
NEAR_BLACK = RGBColor(0x1A, 0x1A, 0x2E)
GRAY       = RGBColor(0x55, 0x55, 0x66)

IMG = '/home/user/Aeropython'

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

FOOTER = 'QCC Tech Fee Program  |  Mentors: Armendariz & Cheung  |  Keywords: air flow · vortex model · weather forecast'


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide():
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, text, l, t, w, h, size=18, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, wrap=True, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def img(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        s = slide.shapes.add_shape(1, l, t, w, h)
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        s.line.color.rgb = RGBColor(0xAA, 0xAA, 0xAA)
        tf = s.text_frame
        tf.text = f'[{os.path.basename(path)}]'


def header(slide, title):
    rect(slide, 0, 0, W, Inches(1.05), DARK_BLUE)
    rect(slide, 0, 0, Inches(0.12), Inches(1.05), ACCENT)
    txt(slide, title, Inches(0.25), Inches(0.15), Inches(12.8), Inches(0.75),
        size=30, bold=True, color=WHITE)


def footer_bar(slide):
    rect(slide, 0, H - Inches(0.38), W, Inches(0.38), MID_BLUE)
    txt(slide, FOOTER, Inches(0.2), H - Inches(0.36), W - Inches(0.4), Inches(0.32),
        size=8, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


def bullets(slide, items, l, t, w, gap=0.52, base_size=17):
    cy = t
    for item in items:
        is_sub = item.startswith('  ')
        text = ('• ' if not is_sub else '  – ') + item.strip()
        size = base_size if not is_sub else base_size - 2
        color = NEAR_BLACK if not is_sub else GRAY
        txt(slide, text, l, cy, w, Inches(0.5), size=size, color=color)
        cy += Inches(gap if not is_sub else gap - 0.06)
    return cy


def two_col(slide, left_items, img_path, img_scale=1.0):
    """Bullets on left, image on right."""
    header_h = Inches(1.1)
    footer_h = Inches(0.4)
    usable_h = H - header_h - footer_h
    col_w = Inches(6.0)
    img_w = Inches(6.8) * img_scale
    img_h = usable_h * img_scale
    img_l = W - img_w - Inches(0.15)
    img_t = header_h + (usable_h - img_h) / 2
    img(slide, img_path, img_l, img_t, img_w, img_h)
    bullets(slide, left_items, Inches(0.25), header_h + Inches(0.2), col_w)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, DARK_BLUE)
rect(s, 0, Inches(2.2), W, Inches(0.08), ACCENT)
rect(s, 0, Inches(5.3), W, Inches(0.08), ACCENT)

txt(s, 'Vortex Modeling with Application to\nAtmospheric Air Flow Pattern Data Analysis',
    Inches(0.6), Inches(2.35), Inches(12.1), Inches(2.8),
    size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, 'QCC Tech Fee Program',
    Inches(0.6), Inches(5.5), Inches(12.1), Inches(0.45),
    size=17, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

txt(s, 'Faculty Mentors: Armendariz & Cheung',
    Inches(0.6), Inches(5.95), Inches(12.1), Inches(0.45),
    size=17, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)

txt(s, 'Keywords: air flow  ·  vortex model  ·  weather forecast',
    Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.4),
    size=13, color=RGBColor(0x88, 0xAA, 0xDD), align=PP_ALIGN.CENTER, italic=True)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — What is a Vortex?
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'What is a Vortex?')
footer_bar(s)

two_col(s, [
    'A vortex is a region of fluid (air or water) rotating around a central axis',
    'The rotating core is called the vortex center or eye',
    'Fluid spins faster near the center, slower at the edges',
    'Vortices are everywhere in nature:',
    '  Hurricanes & typhoons',
    '  Tornadoes',
    '  Ocean gyres',
    '  Aircraft wake turbulence',
    '  Bathtub drain swirls',
    'Understanding vortices is key to weather prediction and aviation safety',
], f'{IMG}/vortex_g5.0.png')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Real-World Connection
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Vortices in Our Atmosphere — Ventusky')
footer_bar(s)

two_col(s, [
    'Ventusky is a real-time wind & weather visualization tool',
    'It shows actual atmospheric vortex patterns from global forecast models',
    'You can see multiple interacting vortices at once',
    'Our simulations reproduce this kind of multi-vortex structure from first principles',
    'Color = wind speed  |  Streamlines = flow direction',
    '  Red/warm = fast',
    '  Blue/cool = slow',
    'Same jet colormap used in our Python simulations',
], f'{IMG}/ventusky.png')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Why Study Vortices? + Motivation
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Why Study Vortices?  |  Project Motivation')
footer_bar(s)

bullets(s, [
    'Weather forecasting accuracy depends on understanding vortex dynamics',
    'Aircraft wake vortices are a major aviation safety concern',
    'Ocean vortices (gyres) control heat transport and climate',
    'Vortex models are used in computational fluid dynamics (CFD) for engineering',
    '',
    'This Project:',
    '  Started from open-source Python vortex codes on GitHub',
    '  Extended with new input conditions: oval shapes, viscosity, multi-vortex systems',
    '  Applied Lamb-Oseen viscous decay — the standard model for real vortices',
    '  Demonstrated a proof-of-concept for atmospheric air flow data analysis',
    '  Developed as a lab exercise tool for community college students',
], Inches(0.4), Inches(1.15), Inches(12.5), gap=0.48)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — The Physics: Potential Flow
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'The Physics — Potential Flow Theory')
footer_bar(s)

bullets(s, [
    'Potential flow: idealized, inviscid (frictionless), irrotational fluid motion',
    'Fluid velocity is derived from a scalar potential or stream function ψ',
    'Key property: flows can be superposed (added together linearly)',
    '  One vortex + another vortex = combined flow field',
    '  Vortex + sink = spiral inward flow (like a drain)',
    '',
    'Two fundamental building blocks used in this project:',
    '  Point Vortex — pure rotation around a fixed center',
    '  Sink — fluid drawn inward from all directions',
    '',
    'Real vortices also have viscosity — handled by the Lamb-Oseen model',
], Inches(0.4), Inches(1.15), Inches(12.5), gap=0.48)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Key Formulas
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Key Formulas')
footer_bar(s)

# Formula boxes
def formula_box(slide, label, formula, l, t, w, h):
    rect(slide, l, t, w, h, LIGHT_BLUE)
    strip = slide.shapes.add_shape(1, l, t, Inches(0.08), h)
    strip.fill.solid(); strip.fill.fore_color.rgb = MID_BLUE; strip.line.fill.background()
    txt(slide, label, l + Inches(0.15), t + Inches(0.05), w - Inches(0.2), Inches(0.3),
        size=13, bold=True, color=MID_BLUE)
    txt(slide, formula, l + Inches(0.15), t + Inches(0.35), w - Inches(0.2), h - Inches(0.4),
        size=15, color=NEAR_BLACK, bold=False)

formula_box(s, 'Stream Function (Point Vortex)',
            'ψ = Γ / (4π) · ln(r²)',
            Inches(0.4), Inches(1.2), Inches(5.9), Inches(0.85))

formula_box(s, 'Velocity Field',
            'u = +Γ/(2π) · (y − y₀) / r²\nv = −Γ/(2π) · (x − x₀) / r²',
            Inches(6.8), Inches(1.2), Inches(6.2), Inches(0.95))

formula_box(s, 'Elliptical (Oval) Distance',
            'r² = (Δx / aₓ)² + (Δy / aᵧ)²\naₓ = aspect_x,  aᵧ = aspect_y',
            Inches(0.4), Inches(2.3), Inches(5.9), Inches(0.95))

formula_box(s, 'Lamb-Oseen Core Radius',
            'r_c(t) = √(4νt + r₀²)\nν = kinematic viscosity,  r₀ = initial core size',
            Inches(6.8), Inches(2.4), Inches(6.2), Inches(0.95))

formula_box(s, 'Lamb-Oseen Velocity (with viscous cutoff)',
            'u = +Γ/(2π) · (Δy/aᵧ²) / r² · [1 − exp(−r²/r_c²)]\nv = −Γ/(2π) · (Δx/aₓ²) / r² · [1 − exp(−r²/r_c²)]',
            Inches(0.4), Inches(3.45), Inches(12.6), Inches(1.05))

txt(s, 'Γ = circulation (strength)   |   r = distance from center   |   ν = viscosity   |   t = time',
    Inches(0.4), Inches(4.7), Inches(12.5), Inches(0.4),
    size=13, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

txt(s, 'The [1 − exp(…)] term removes the singularity at r = 0, making the vortex core smooth and physically realistic.',
    Inches(0.4), Inches(5.2), Inches(12.5), Inches(0.5),
    size=14, color=NEAR_BLACK, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Step 1: Basic Vortex (vortex.py)
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Step 1 — Basic Point Vortex  (vortex.py)')
footer_bar(s)

two_col(s, [
    'Starting point: a single ideal point vortex',
    'Circular streamlines — fluid orbits the center',
    'Color = normalized flow speed',
    '  Red/warm near center = fast',
    '  Blue/cool at edges = slow',
    'No viscosity — flow speed → ∞ at center (singularity)',
    'GitHub source code adapted and verified',
    '',
    'Γ = 5.0  |  Center: (0.0, −0.3)',
], f'{IMG}/vortex_g5.0.png')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Step 2: Vortex Movement (calculate_vortex.py)
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Step 2 — Simulating Vortex Movement  (calculate_vortex.py)')
footer_bar(s)

two_col(s, [
    'Extended the base code to track the trajectory of the vortex center',
    'Vortex moves through the domain over time',
    'Dots show sampled positions along the orbit path',
    'Red dot = vortex center  |  Blue dot = reference point',
    '',
    'Demonstrates that even a single vortex in a bounded domain',
    'produces non-trivial motion paths',
    '',
    'Foundation for the N-vortex mutual interaction model',
], f'{IMG}/calculate.png')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Step 3: Vortex + Sink
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Step 3 — Superposition: Vortex + Sink  (vortex.py)')
footer_bar(s)

two_col(s, [
    'Added a sink to the vortex — fluid spirals inward',
    'Superposition principle: velocity fields add linearly',
    'Result: spiral streamlines (like a draining bathtub)',
    '',
    'Physically models:',
    '  Low-pressure systems pulling air inward while rotating',
    '  Hurricane eye-wall dynamics',
    '  Drain vortices',
    '',
    'Γ = 5.0  |  Sink strength = 20.0',
    'Sink at (0.25, −0.5)',
], f'{IMG}/vortex_g5.0_s20.0_(0.25,-0.5).png')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Lamb-Oseen Theory
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Lamb-Oseen Viscous Decay Model')
footer_bar(s)

bullets(s, [
    'Real vortices are viscous — they diffuse and weaken over time',
    'The Lamb-Oseen model is the standard physics model for this effect',
    '',
    'Core radius grows with time:   r_c(t) = √(4νt + r₀²)',
    '  At t = 0: core = r₀ (initial size)',
    '  As t increases: core spreads due to viscous diffusion',
    '  Larger ν = faster spreading',
    '',
    'The [1 − exp(−r²/r_c²)] factor:',
    '  = 0 at the center → no singularity',
    '  → 1 far from center → recovers ideal vortex behavior',
    '',
    'In this project: ν = 0.005  (scaled for visible animation)',
    'Simulation runs from t = 0.5 to t = 12.5  (120 frames, Δt = 0.1)',
], Inches(0.4), Inches(1.15), Inches(12.5), gap=0.45)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Step 4: Lamb-Oseen Simulation
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Step 4 — Multi-Vortex Lamb-Oseen Simulation  (lamb_oseen.py)')
footer_bar(s)

two_col(s, [
    '4 vortices with different strengths and positions',
    'Red dots = counter-clockwise (Γ > 0)',
    'Blue dots = clockwise (Γ < 0)',
    'Dashed ellipses = viscous core boundary (grows over time)',
    'Particles (gray dots) advected by the combined flow',
    '',
    'Oval shape: each vortex has aspect_x / aspect_y',
    '  aspect_x > 1 → wider horizontally',
    '  aspect_y > 1 → taller vertically',
    '  Mimics real elongated atmospheric vortices',
    '',
    'Animation shows core growth as viscosity diffuses the vortex',
], f'{IMG}/lamb_oseen.png')


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Oval Vortex Shape
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Oval Vortex Cores — Aspect Ratio Parameters')
footer_bar(s)

bullets(s, [
    'Real atmospheric vortices are rarely perfectly circular',
    'Elongated by wind shear, terrain, pressure gradients',
    '',
    'Implemented via elliptical distance metric:',
    '  r² = (Δx / aspect_x)² + (Δy / aspect_y)²',
    '',
    'Vortex parameters used in the simulation:',
], Inches(0.4), Inches(1.15), Inches(12.5), gap=0.5)

# Table
col_labels = ['Vortex', 'Strength (Γ)', 'aspect_x', 'aspect_y', 'Shape']
rows = [
    ['1', '−5.0 (CW)', '1.5', '0.7', 'Wide oval'],
    ['2', '−4.0 (CW)', '1.0', '1.0', 'Circular'],
    ['3', '+5.0 (CCW)', '0.6', '1.4', 'Tall oval'],
    ['4', '−1.0 (CW)', '1.2', '0.8', 'Slight wide oval'],
]
col_w = [Inches(1.2), Inches(2.2), Inches(1.6), Inches(1.6), Inches(2.2)]
col_x = [Inches(0.5 + sum(col_w[:i].count(col_w[i]) for i in range(0))),
         Inches(1.7), Inches(3.9), Inches(5.5), Inches(7.1)]
# recompute
cx = Inches(0.5)
col_positions = []
for cw in col_w:
    col_positions.append(cx)
    cx += cw

row_h = Inches(0.42)
t0 = Inches(4.3)

for ci, (label, cp) in enumerate(zip(col_labels, col_positions)):
    rect(s, cp, t0, col_w[ci], row_h, DARK_BLUE)
    txt(s, label, cp + Inches(0.05), t0 + Inches(0.06), col_w[ci] - Inches(0.1), row_h,
        size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(rows):
    bg_color = LIGHT_BLUE if ri % 2 == 0 else WHITE
    for ci, (cell, cp) in enumerate(zip(row, col_positions)):
        rect(s, cp, t0 + row_h * (ri + 1), col_w[ci], row_h, bg_color)
        txt(s, cell, cp + Inches(0.05), t0 + row_h * (ri + 1) + Inches(0.06),
            col_w[ci] - Inches(0.1), row_h, size=13, color=NEAR_BLACK, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — CV = a/b
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Coefficient of Variation as Diversity Measure  (CV = a/b)')
footer_bar(s)

bullets(s, [
    'To characterize how elongated each vortex core is, we define:',
    '',
], Inches(0.4), Inches(1.15), Inches(12.5), gap=0.4)

formula_box(s, 'Semi-axes of the viscous core ellipse',
            'a = r_c(t) · aspect_x        (x-radius)\nb = r_c(t) · aspect_y        (y-radius)',
            Inches(0.8), Inches(2.3), Inches(5.6), Inches(1.05))

formula_box(s, 'Coefficient of Variation',
            'CV = a / b = aspect_x / aspect_y',
            Inches(7.0), Inches(2.3), Inches(5.9), Inches(1.05))

txt(s, 'Note: since both a and b scale with the same r_c(t), CV is constant over time for each vortex.',
    Inches(0.4), Inches(3.55), Inches(12.5), Inches(0.45),
    size=14, color=GRAY, italic=True, align=PP_ALIGN.CENTER)

bullets(s, [
    'CV = 1.0  →  perfectly circular core',
    'CV > 1.0  →  core is wider than tall  (stretched horizontally)',
    'CV < 1.0  →  core is taller than wide  (stretched vertically)',
    'Used as a diversity measure to compare shape configurations across vortices',
    'Supports weather forecast research by quantifying vortex shape asymmetry',
], Inches(0.4), Inches(4.1), Inches(12.5), gap=0.48)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — CV Results
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'CV Results — Vortex Shape Diversity')
footer_bar(s)

txt(s, 'Computed at t = 0.5  (initial frame)  via print_cv() in lamb_oseen.py:',
    Inches(0.4), Inches(1.2), Inches(12.5), Inches(0.4), size=15, color=NEAR_BLACK)

cv_cols = ['Vortex', 'aspect_x', 'aspect_y', 'r_c (t=0.5)', 'a = r_c·aspect_x', 'b = r_c·aspect_y', 'CV = a/b']
cv_data = [
    ['1', '1.5', '0.7', '0.1003', '0.1505', '0.0702', '2.1429'],
    ['2', '1.0', '1.0', '0.1003', '0.1003', '0.1003', '1.0000'],
    ['3', '0.6', '1.4', '0.1049', '0.0629', '0.1469', '0.4286'],
    ['4', '1.2', '0.8', '0.1000', '0.1200', '0.0800', '1.5000'],
]
cv_col_w = [Inches(1.1), Inches(1.3), Inches(1.3), Inches(1.7), Inches(2.2), Inches(2.2), Inches(1.7)]
cx2 = Inches(0.3)
cv_positions = []
for cw in cv_col_w:
    cv_positions.append(cx2)
    cx2 += cw

t0 = Inches(1.8)
rh = Inches(0.44)

for ci, (label, cp) in enumerate(zip(cv_cols, cv_positions)):
    rect(s, cp, t0, cv_col_w[ci], rh, DARK_BLUE)
    txt(s, label, cp + Inches(0.03), t0 + Inches(0.06), cv_col_w[ci] - Inches(0.06), rh,
        size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

for ri, row in enumerate(cv_data):
    bg_c = LIGHT_BLUE if ri % 2 == 0 else WHITE
    for ci, (cell, cp) in enumerate(zip(row, cv_positions)):
        is_cv = (ci == 6)
        rect(s, cp, t0 + rh * (ri + 1), cv_col_w[ci], rh, bg_c if not is_cv else RGBColor(0xFF, 0xF0, 0xD0))
        txt(s, cell, cp + Inches(0.03), t0 + rh * (ri + 1) + Inches(0.07),
            cv_col_w[ci] - Inches(0.06), rh, size=12, color=NEAR_BLACK if not is_cv else ACCENT,
            bold=is_cv, align=PP_ALIGN.CENTER)

bullets(s, [
    'Vortex 1 (CV=2.14): most elongated horizontally — mimics large-scale weather fronts',
    'Vortex 2 (CV=1.00): circular — baseline reference',
    'Vortex 3 (CV=0.43): most elongated vertically — like a narrow storm column',
    'Vortex 4 (CV=1.50): moderately wide — typical mid-latitude cyclone shape',
], Inches(0.4), Inches(4.2), Inches(12.5), gap=0.50)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Atmospheric Application
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Application to Atmospheric Air Flow Data — Proof of Concept')
footer_bar(s)

two_col(s, [
    'Ventusky shows real atmospheric vortex patterns from global weather models',
    'Our model reproduces key features:',
    '  Multiple interacting vortex centers',
    '  Varying rotation directions (CW / CCW)',
    '  Non-circular, asymmetric cores',
    '  Speed gradient: fast core, slow periphery',
    '',
    'CV as a metric could be used to:',
    '  Compare simulated vs observed vortex shapes',
    '  Track how vortex geometry changes over time',
    '  Classify storm types by shape diversity',
    '',
    'Next step: overlay model output on real forecast data',
], f'{IMG}/ventusky.png', img_scale=0.92)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — Cyberlearning Extension
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Cyberlearning Lab Exercise Extension')
footer_bar(s)

bullets(s, [
    'Goal: develop an interactive vortex lab for community college students',
    '',
    'Students can:',
    '  Modify vortex parameters (strength, position, aspect ratio, viscosity)',
    '  Run simulations directly in a web browser (no install needed)',
    '  Compare model output to Ventusky real-world atmospheric data',
    '  Calculate and interpret CV values as a shape diversity measure',
    '',
    'Tools:',
    '  Python + Jupyter Notebook (JupyterHub / GitHub Codespaces)',
    '  matplotlib animations rendered in-browser',
    '  GitHub repository for version control and collaboration',
    '',
    'Pedagogical value: connects fluid physics, coding, and real-world weather data',
    'Accessible to students with no prior programming experience',
], Inches(0.4), Inches(1.15), Inches(12.5), gap=0.45)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Conclusions + Acknowledgements
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
bg(s, WHITE)
header(s, 'Conclusions & Acknowledgements')
footer_bar(s)

txt(s, 'Conclusions', Inches(0.4), Inches(1.15), Inches(12.5), Inches(0.4),
    size=20, bold=True, color=DARK_BLUE)
bullets(s, [
    'Successfully extended open-source Python vortex codes for multiple input conditions',
    'Implemented oval/asymmetric vortex cores via elliptical distance metric',
    'Applied Lamb-Oseen viscous decay — physically realistic vortex model',
    'CV = a/b provides a simple, interpretable shape diversity measure',
    'Demonstrated atmospheric air flow pattern analysis as a proof of concept',
    'Framework ready for interactive cyberlearning lab development',
], Inches(0.4), Inches(1.6), Inches(12.5), gap=0.44)

txt(s, 'Acknowledgements', Inches(0.4), Inches(4.85), Inches(12.5), Inches(0.4),
    size=20, bold=True, color=DARK_BLUE)

rect(s, Inches(0.4), Inches(5.3), Inches(12.5), Inches(1.5), LIGHT_BLUE)
txt(s, 'This project was supported by the QCC Tech Fee Program.\nFaculty Mentors: Armendariz and Cheung\nCommunity college students participated in weather forecast research using cyberlearning tools.',
    Inches(0.55), Inches(5.38), Inches(12.2), Inches(1.3),
    size=15, color=DARK_BLUE, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════════════════
out = '/home/user/Aeropython/vortex_presentation.pptx'
prs.save(out)
print(f'Saved: {out}')
